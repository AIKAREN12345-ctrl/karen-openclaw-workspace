import pandas as pd
import matplotlib.pyplot as plt
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
import io

# Read the data
df = pd.read_excel('C:/Users/Karen/.openclaw/media/inbound/file_191---1d6a3fcf-e32b-413e-8746-b239ed12319a.xlsx')
df = df[df['Sales Person'] != 'Sales Person']
df['Date'] = pd.to_datetime(df['Date'])
df['Amount'] = pd.to_numeric(df['Amount'], errors='coerce')
df['Boxes Shipped'] = pd.to_numeric(df['Boxes Shipped'], errors='coerce')

# Calculate metrics
total_revenue = df['Amount'].sum()
total_transactions = len(df)
unique_products = df['Product'].nunique()
unique_countries = df['Country'].nunique()
avg_deal_size = total_revenue / total_transactions

# Create presentation - 16:9 format
prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

# Color scheme
DARK_BLUE = RGBColor(31, 78, 120)
LIGHT_BLUE = RGBColor(68, 114, 196)
WHITE = RGBColor(255, 255, 255)
GRAY = RGBColor(128, 128, 128)

def add_title_slide(prs, title, subtitle):
    blank_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(blank_layout)
    
    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height)
    bg.fill.solid()
    bg.fill.fore_color.rgb = DARK_BLUE
    bg.line.fill.background()
    
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(2.5), Inches(12.333), Inches(1.5))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(44)
    p.font.bold = True
    p.font.color.rgb = WHITE
    p.alignment = PP_ALIGN.CENTER
    
    sub_box = slide.shapes.add_textbox(Inches(0.5), Inches(4.2), Inches(12.333), Inches(1))
    tf = sub_box.text_frame
    p = tf.paragraphs[0]
    p.text = subtitle
    p.font.size = Pt(24)
    p.font.color.rgb = RGBColor(200, 200, 200)
    p.alignment = PP_ALIGN.CENTER
    
    return slide

def add_content_slide(prs, title, content_lines):
    blank_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(blank_layout)
    
    header = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, Inches(1.2))
    header.fill.solid()
    header.fill.fore_color.rgb = DARK_BLUE
    header.line.fill.background()
    
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.25), Inches(12.333), Inches(0.8))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(32)
    p.font.bold = True
    p.font.color.rgb = WHITE
    
    content_box = slide.shapes.add_textbox(Inches(0.7), Inches(1.6), Inches(12), Inches(5.5))
    tf = content_box.text_frame
    tf.word_wrap = True
    
    for i, line in enumerate(content_lines):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = line
        p.font.size = Pt(18)
        p.font.color.rgb = RGBColor(64, 64, 64)
        p.space_after = Pt(12)
    
    return slide

def create_chart_image(chart_type, data, title, xlabel='', ylabel=''):
    fig, ax = plt.subplots(figsize=(10, 5.5))
    
    if chart_type == 'bar':
        bars = ax.bar(data['x'], data['y'], color='#4472C4', edgecolor='#1F4E78', linewidth=1.5)
        ax.set_xticklabels(data['x'], rotation=45, ha='right')
        for bar in bars:
            height = bar.get_height()
            ax.text(bar.get_x() + bar.get_width()/2., height,
                   f'€{height/1000000:.2f}M' if height > 1000000 else f'€{height/1000:.0f}K',
                   ha='center', va='bottom', fontsize=9)
    elif chart_type == 'line':
        ax.plot(data['x'], data['y'], marker='o', linewidth=3, color='#4472C4', markersize=8)
        ax.fill_between(range(len(data['x'])), data['y'], alpha=0.3, color='#4472C4')
        ax.set_xticks(range(len(data['x'])))
        ax.set_xticklabels(data['x'], rotation=45, ha='right')
    
    ax.set_title(title, fontsize=14, fontweight='bold', color='#1F4E78')
    ax.set_xlabel(xlabel, fontsize=11)
    ax.set_ylabel(ylabel, fontsize=11)
    ax.spines['top'].set_visible(False)
    ax.spines['right'].set_visible(False)
    ax.grid(axis='y', alpha=0.3)
    
    plt.tight_layout()
    img_buffer = io.BytesIO()
    plt.savefig(img_buffer, format='png', dpi=150, bbox_inches='tight')
    img_buffer.seek(0)
    plt.close()
    return img_buffer

def add_chart_slide(prs, title, chart_buffer):
    blank_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(blank_layout)
    
    header = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, Inches(1.2))
    header.fill.solid()
    header.fill.fore_color.rgb = DARK_BLUE
    header.line.fill.background()
    
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.25), Inches(12.333), Inches(0.8))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(32)
    p.font.bold = True
    p.font.color.rgb = WHITE
    
    slide.shapes.add_picture(chart_buffer, Inches(1.5), Inches(1.5), width=Inches(10.333))
    
    return slide

# SLIDE 1: Title
add_title_slide(prs, 'Chocolate Sales Analysis', 'Q1-Q4 Performance Review | Data-Driven Insights')

# SLIDE 2: Executive Summary (merged with Key Findings)
exec_summary = [
    f'• Total Revenue: €{total_revenue/1000000:.2f}M across {total_transactions:,} transactions',
    f'• Product Portfolio: {unique_products} unique products in {unique_countries} countries',
    f'• Average Deal Size: €{avg_deal_size:,.0f} per transaction',
    '',
    'Key Insights:',
    '• Australia leads with €2.24M (37.5% of total revenue)',
    '• Smooth Silky Salty Bar is top product at €1.37M',
    '• Peak sales in August (€1.06M); lowest in February (€85K)',
    '• Ches Bonnell is top performer with €1.82M in sales'
]
add_content_slide(prs, 'Executive Summary', exec_summary)

# SLIDE 3: Geographic Analysis
revenue_by_country = df.groupby('Country')['Amount'].sum().sort_values(ascending=False)
country_data = {'x': list(revenue_by_country.index), 'y': list(revenue_by_country.values)}
chart_buf = create_chart_image('bar', country_data, '', '', 'Revenue (€)')
add_chart_slide(prs, 'Revenue by Country', chart_buf)

# SLIDE 4: Product Performance
top_products = df.groupby('Product')['Amount'].sum().sort_values(ascending=False).head(5)
product_data = {'x': list(top_products.index), 'y': list(top_products.values)}
chart_buf = create_chart_image('bar', product_data, '', '', 'Revenue (€)')
add_chart_slide(prs, 'Top 5 Products by Revenue', chart_buf)

# SLIDE 5: Monthly Trends
monthly_revenue = df.groupby(df['Date'].dt.to_period('M'))['Amount'].sum()
month_labels = [str(m) for m in monthly_revenue.index]
month_data = {'x': month_labels, 'y': list(monthly_revenue.values)}
chart_buf = create_chart_image('line', month_data, '', 'Month', 'Revenue (€)')
add_chart_slide(prs, 'Monthly Revenue Trends', chart_buf)

# SLIDE 6: Sales Team Performance
team_performance = df.groupby('Sales Person')['Amount'].sum().sort_values(ascending=False)
team_data = {'x': list(team_performance.index), 'y': list(team_performance.values)}
chart_buf = create_chart_image('bar', team_data, '', '', 'Revenue (€)')
add_chart_slide(prs, 'Sales Team Performance', chart_buf)

# SLIDE 7: Recommendations
recommendations = [
    '1. Expand Australian Market Leadership',
    '   • Australia generates 37.5% of revenue—invest in market expansion',
    '   • Replicate successful strategies in other regions',
    '',
    '2. Address Seasonal Volatility',
    '   • August peak (€1.06M) vs February low (€85K) shows 12x variation',
    '   • Develop off-season promotions to smooth revenue',
    '',
    '3. Leverage Top Performers',
    '   • Ches Bonnell (€1.82M) and Gigi Bohling (€1.42M) lead the team',
    '   • Share best practices and consider mentorship programs',
    '',
    '4. Product Portfolio Optimization',
    '   • Smooth Silky Salty Bar dominates—ensure supply chain resilience',
    '   • Consider expanding premium product lines'
]
add_content_slide(prs, 'Strategic Recommendations', recommendations)

# SLIDE 8: Thank You
add_title_slide(prs, 'Thank You', 'Questions & Discussion')

# Save presentation
prs.save('C:/Users/Karen/.openclaw/workspace/Chocolate_Sales_Presentation.pptx')
print('SUCCESS: Presentation created: 8 slides (down from 12)')
print('Slides: Title -> Executive Summary -> Country -> Products -> Trends -> Team -> Recommendations -> Thank You')
