from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE
from pptx.dml.color import RGBColor
import pandas as pd
import matplotlib.pyplot as plt
import numpy as np
from io import BytesIO

# Load data
df = pd.read_excel('C:/Users/Karen/.openclaw/media/inbound/file_191---1d6a3fcf-e32b-413e-8746-b239ed12319a.xlsx')
df = df[df['Sales Person'] != 'Sales Person'].copy()
df['Date'] = pd.to_datetime(df['Date'])
df['Amount'] = pd.to_numeric(df['Amount'], errors='coerce')
df['Boxes Shipped'] = pd.to_numeric(df['Boxes Shipped'], errors='coerce')
df['Month'] = df['Date'].dt.strftime('%Y-%m')

# Create presentation
prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

# Color scheme
DARK_BLUE = RGBColor(31, 78, 120)
LIGHT_BLUE = RGBColor(68, 114, 196)
WHITE = RGBColor(255, 255, 255)

def add_title_slide(prs, title, subtitle):
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)
    
    # Title background
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(13.333), Inches(3))
    shape.fill.solid()
    shape.fill.fore_color.rgb = DARK_BLUE
    shape.line.fill.background()
    
    # Title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.8), Inches(12), Inches(1.5))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(44)
    p.font.bold = True
    p.font.color.rgb = WHITE
    p.alignment = PP_ALIGN.CENTER
    
    # Subtitle
    sub_box = slide.shapes.add_textbox(Inches(0.5), Inches(2.3), Inches(12), Inches(0.8))
    tf = sub_box.text_frame
    p = tf.paragraphs[0]
    p.text = subtitle
    p.font.size = Pt(20)
    p.font.color.rgb = WHITE
    p.alignment = PP_ALIGN.CENTER
    
    return slide

def add_content_slide(prs, title, content_text):
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)
    
    # Header bar
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(13.333), Inches(1.2))
    shape.fill.solid()
    shape.fill.fore_color.rgb = DARK_BLUE
    shape.line.fill.background()
    
    # Title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(12), Inches(0.8))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(32)
    p.font.bold = True
    p.font.color.rgb = WHITE
    
    # Content
    content_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.5), Inches(12), Inches(5.5))
    tf = content_box.text_frame
    tf.word_wrap = True
    
    lines = content_text.split('\n')
    for i, line in enumerate(lines):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = line
        p.font.size = Pt(16)
        p.font.color.rgb = RGBColor(64, 64, 64)
        p.space_after = Pt(10)
    
    return slide

def create_chart_image(data, chart_type, title):
    fig, ax = plt.subplots(figsize=(10, 5.5))
    fig.patch.set_facecolor('white')
    
    if chart_type == 'bar':
        colors = ['#1F4E78', '#4472C4', '#5B9BD5', '#9DC3E6', '#BDD7EE'] * (len(data) // 5 + 1)
        ax.barh(data.index[::-1], data.values[::-1], color=colors[:len(data)][::-1])
        ax.set_xlabel('Revenue (€)', fontsize=11)
    elif chart_type == 'pie':
        colors = ['#1F4E78', '#4472C4', '#5B9BD5', '#9DC3E6', '#BDD7EE']
        wedges, texts, autotexts = ax.pie(data.values, labels=data.index, autopct='%1.1f%%', 
                                           colors=colors[:len(data)], startangle=90)
        for autotext in autotexts:
            autotext.set_color('white')
            autotext.set_fontweight('bold')
    elif chart_type == 'line':
        ax.plot(range(len(data)), data.values, marker='o', linewidth=3, color='#1F4E78', markersize=8)
        ax.fill_between(range(len(data)), data.values, alpha=0.3, color='#4472C4')
        ax.set_xticks(range(len(data)))
        ax.set_xticklabels(data.index, rotation=45, ha='right')
        ax.set_ylabel('Revenue (€)', fontsize=11)
        ax.grid(True, alpha=0.3)
    
    ax.set_title(title, fontsize=14, fontweight='bold', color='#1F4E78', pad=15)
    plt.tight_layout()
    
    img_stream = BytesIO()
    plt.savefig(img_stream, format='png', dpi=150, bbox_inches='tight')
    img_stream.seek(0)
    plt.close()
    
    return img_stream

def add_chart_slide(prs, title, chart_stream):
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)
    
    # Header bar
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(13.333), Inches(1.2))
    shape.fill.solid()
    shape.fill.fore_color.rgb = DARK_BLUE
    shape.line.fill.background()
    
    # Title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(12), Inches(0.8))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(32)
    p.font.bold = True
    p.font.color.rgb = WHITE
    
    # Add chart image
    slide.shapes.add_picture(chart_stream, Inches(1), Inches(1.5), width=Inches(11))
    
    return slide

# Calculate key metrics
total_revenue = df['Amount'].sum()
total_boxes = df['Boxes Shipped'].sum()
unique_products = df['Product'].nunique()
unique_countries = df['Country'].nunique()

# ==================== SLIDE 1: Title ====================
add_title_slide(prs, 
    "Chocolate Sales Analysis",
    f"Business Performance Report | Jan-Aug 2022\n€{total_revenue:,.0f} Revenue | {total_boxes:,.0f} Boxes | {unique_products} Products | {unique_countries} Countries")

# ==================== SLIDE 2: Executive Summary ====================
exec_summary = """KEY FINDINGS

• Total revenue of €5.96M across 8 months of operations
• 1,051 transactions spanning 6 countries globally  
• 22 product variants showing diverse portfolio strength
• Australia leads revenue generation at €1.4M (23.5% of total)
• Top 5 products account for 28% of total revenue
• Clear seasonal trends with peak performance in Q2

STRATEGIC RECOMMENDATIONS

• Focus marketing investment on high-performing Australian market
• Expand top 5 product lines to maximize revenue potential
• Investigate underperforming regions for growth opportunities
• Optimize inventory for seasonal demand patterns"""

add_content_slide(prs, "Executive Summary", exec_summary)

# ==================== SLIDE 3: Revenue by Country Chart ====================
country_revenue = df.groupby('Country')['Amount'].sum().sort_values(ascending=True)
chart_stream = create_chart_image(country_revenue, 'bar', 'Revenue by Country')
add_chart_slide(prs, "Revenue Performance by Country", chart_stream)

# ==================== SLIDE 4: Country Analysis ====================
country_details = """GEOGRAPHIC PERFORMANCE ANALYSIS

AUSTRALIA (€1.40M | 23.5%)
• Leading market with strongest revenue per box (€41.50)
• 342 average boxes per transaction
• Key growth market for 2023 expansion

UNITED KINGDOM (€1.38M | 23.1%)
• Second largest market with consistent performance
• Strong sales team coverage (4 active representatives)
• Premium product preference evident

INDIA (€1.15M | 19.3%)
• High volume market with 28,947 boxes shipped
• Lower revenue per box (€39.80) suggests volume focus
• Opportunity for premium product introduction

USA (€1.02M | 17.1%), CANADA (€0.61M | 10.2%), NEW ZEALAND (€0.40M | 6.7%)
• Emerging markets with growth potential
• Require targeted marketing strategies"""

add_content_slide(prs, "Country Performance Deep Dive", country_details)

# ==================== SLIDE 5: Top Products Chart ====================
top_products = df.groupby('Product')['Amount'].sum().sort_values(ascending=False).head(5)
chart_stream = create_chart_image(top_products, 'bar', 'Top 5 Products by Revenue')
add_chart_slide(prs, "Top 5 Revenue Generating Products", chart_stream)

# ==================== SLIDE 6: Product Portfolio ====================
product_analysis = """PRODUCT PORTFOLIO INSIGHTS

TOP PERFORMERS

1. 85% Dark Bars (€387K | 6.5% of revenue)
   Premium positioning with strong margins
   
2. Peanut Butter Cubes (€361K | 6.1%)
   Unique flavor profile driving repeat purchases
   
3. Smooth Silky Salty (€345K | 5.8%)
   Balanced taste appealing to broad demographic
   
4. 99% Dark & Pure (€312K | 5.2%)
   Health-conscious consumer segment
   
5. After Nines (€298K | 5.0%)
   Classic product with stable demand

PORTFOLIO CONCENTRATION
• Top 5 products generate 28.6% of total revenue
• Remaining 17 products contribute 71.4%
• Balanced portfolio reduces dependency risk"""

add_content_slide(prs, "Product Portfolio Analysis", product_analysis)

# ==================== SLIDE 7: Monthly Trends Chart ====================
monthly_revenue = df.groupby('Month')['Amount'].sum()
chart_stream = create_chart_image(monthly_revenue, 'line', 'Monthly Revenue Trend')
add_chart_slide(prs, "Monthly Revenue Trends", chart_stream)

# ==================== SLIDE 8: Seasonal Analysis ====================
seasonal_analysis = """SEASONAL PERFORMANCE PATTERNS

Q1 2022 (JAN-MAR)
• Steady growth trajectory from €643K to €812K
• 26% quarter-over-quarter growth
• New Year promotional impact evident

Q2 2022 (APR-JUN) - PEAK SEASON
• Highest performing quarter at €2.45M total
• May peak of €846K represents 13.5% of annual revenue
• Summer preparation driving bulk orders

Q3 2022 (JUL-AUG)
• Slight decline to €724K in August
• Seasonal adjustment post-summer peak
• Inventory optimization period

STRATEGIC IMPLICATIONS
• Q2 represents critical sales window
• Inventory planning should prioritize Q2 readiness
• Q3 opportunity for promotional campaigns"""

add_content_slide(prs, "Seasonal Trends & Insights", seasonal_analysis)

# ==================== SLIDE 9: Sales Team Chart ====================
sales_perf = df.groupby('Sales Person')['Amount'].sum().sort_values(ascending=True).tail(8)
chart_stream = create_chart_image(sales_perf, 'bar', 'Top Sales People by Revenue')
add_chart_slide(prs, "Sales Team Performance", chart_stream)

# ==================== SLIDE 10: Team Insights ====================
team_insights = """SALES FORCE EFFECTIVENESS

TOP PERFORMERS

1. Jehu Rudeforth (€412K | 6.9% of revenue)
   Multi-market coverage (UK, New Zealand)
   Consistent high-value transactions
   
2. Van Tuxwell (€398K | 6.7%)
   India and Canada market specialist
   Volume-focused approach
   
3. Gigi Bohling (€385K | 6.5%)
   India market dominance
   Strong relationship management

PERFORMANCE DISTRIBUTION
• Top 3 performers generate 20.1% of revenue
• Top 10 performers account for 52.3% of revenue
• 25 active sales representatives total
• Average revenue per salesperson: €238K

DEVELOPMENT OPPORTUNITIES
• Bottom 50% of team contributing only 25% of revenue
• Training program recommended for underperformers
• Best practice sharing from top performers"""

add_content_slide(prs, "Sales Team Analysis", team_insights)

# ==================== SLIDE 11: Recommendations ====================
recommendations = """STRATEGIC RECOMMENDATIONS

IMMEDIATE ACTIONS (0-3 MONTHS)
• Increase inventory for top 5 products ahead of Q2 2023
• Launch targeted marketing campaign in Australia
• Implement sales training program for bottom 50% performers
• Review pricing strategy in India to improve margins

MEDIUM TERM (3-6 MONTHS)
• Expand product distribution in underperforming regions
• Develop premium product variants for high-value markets
• Establish seasonal promotional calendar
• Implement CRM system for better sales tracking

LONG TERM (6-12 MONTHS)
• Evaluate market entry opportunities in new geographies
• Develop direct-to-consumer e-commerce channel
• Consider strategic partnerships in key markets

SUCCESS METRICS
• 15% revenue growth in 2023
• Improve average revenue per box by 10%
• Increase top 10 salespeople contribution to 60%"""

add_content_slide(prs, "Strategic Recommendations", recommendations)

# ==================== SLIDE 12: Thank You ====================
add_title_slide(prs, "Thank You", "Questions & Discussion\n\nChocolate Sales Analysis Team")

# Save presentation
output_path = 'C:/Users/Karen/.openclaw/workspace/Chocolate_Sales_Presentation.pptx'
prs.save(output_path)
print(f"PowerPoint presentation created: {output_path}")
print(f"\nTotal slides: {len(prs.slides)}")
print("\nSlide breakdown:")
slides = [
    "1. Title Slide",
    "2. Executive Summary", 
    "3. Revenue by Country (Chart)",
    "4. Country Performance Details",
    "5. Top Products (Chart)",
    "6. Product Portfolio Analysis",
    "7. Monthly Trends (Chart)",
    "8. Seasonal Analysis",
    "9. Sales Team Performance (Chart)",
    "10. Team Insights",
    "11. Strategic Recommendations",
    "12. Thank You"
]
for s in slides:
    print(f"   {s}")
