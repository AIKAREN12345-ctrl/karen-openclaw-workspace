from pptx import Presentation
import sys

files = [
    'C:/Users/Karen/.openclaw/media/inbound/file_202---fc0f38ab-cfbd-4b21-9717-ef421f7b92e2.pptx',
    'C:/Users/Karen/.openclaw/media/inbound/file_203---b29a81c2-ba25-42fc-9f59-0f0311257355.pptx',
    'C:/Users/Karen/.openclaw/media/inbound/file_204---b447abff-0f41-493c-9c03-fe0cf837a9b0.pptx',
    'C:/Users/Karen/.openclaw/media/inbound/file_205---3ee49f12-da47-45ff-b92b-e9a04c6c9052.pptx',
    'C:/Users/Karen/.openclaw/media/inbound/file_207---0fcc2474-bf41-41b8-86d8-639b0cfb1fb6.pptx'
]

with open('pptx_contents.txt', 'w', encoding='utf-8') as out:
    for i, f in enumerate(files, 1):
        out.write(f'\n{"="*60}\n')
        out.write(f'FILE {i}: {f.split("/")[-1]}\n')
        out.write(f'{"="*60}\n\n')
        try:
            prs = Presentation(f)
            for j, slide in enumerate(prs.slides, 1):
                out.write(f'\n--- Slide {j} ---\n')
                for shape in slide.shapes:
                    if hasattr(shape, 'text') and shape.text.strip():
                        out.write(shape.text + '\n')
        except Exception as e:
            out.write(f'Error: {e}\n')

print('Done reading PowerPoints')