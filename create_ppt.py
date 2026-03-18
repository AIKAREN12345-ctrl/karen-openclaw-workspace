from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

# Paddy Power brand colors
PADDY_GREEN = RGBColor(0, 128, 0)  # Dark green
PADDY_YELLOW = RGBColor(255, 215, 0)  # Gold/yellow
WHITE = RGBColor(255, 255, 255)
BLACK = RGBColor(0, 0, 0)
DARK_GRAY = RGBColor(64, 64, 64)

# Create presentation
prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

def add_title_slide(prs, title, subtitle):
    blank_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(blank_layout)
    
    # Green background
    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height)
    bg.fill.solid()
    bg.fill.fore_color.rgb = PADDY_GREEN
    bg.line.fill.background()
    
    # Yellow accent bar at top
    accent = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, Inches(0.3))
    accent.fill.solid()
    accent.fill.fore_color.rgb = PADDY_YELLOW
    accent.line.fill.background()
    
    # Title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(2.5), Inches(12.333), Inches(1.5))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(44)
    p.font.bold = True
    p.font.color.rgb = WHITE
    p.alignment = PP_ALIGN.CENTER
    
    # Subtitle
    sub_box = slide.shapes.add_textbox(Inches(0.5), Inches(4.2), Inches(12.333), Inches(1))
    tf = sub_box.text_frame
    p = tf.paragraphs[0]
    p.text = subtitle
    p.font.size = Pt(24)
    p.font.color.rgb = PADDY_YELLOW
    p.alignment = PP_ALIGN.CENTER
    
    return slide

def add_content_slide(prs, title, bullets, has_image=False):
    blank_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(blank_layout)
    
    # White background
    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height)
    bg.fill.solid()
    bg.fill.fore_color.rgb = WHITE
    bg.line.fill.background()
    
    # Green header bar
    header = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, Inches(1.2))
    header.fill.solid()
    header.fill.fore_color.rgb = PADDY_GREEN
    header.line.fill.background()
    
    # Yellow accent line
    accent = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, Inches(1.2), prs.slide_width, Inches(0.1))
    accent.fill.solid()
    accent.fill.fore_color.rgb = PADDY_YELLOW
    accent.line.fill.background()
    
    # Title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.25), Inches(12.333), Inches(0.8))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(36)
    p.font.bold = True
    p.font.color.rgb = WHITE
    
    # Content bullets
    content_width = Inches(8) if has_image else Inches(12)
    content_box = slide.shapes.add_textbox(Inches(0.7), Inches(1.6), content_width, Inches(5.5))
    tf = content_box.text_frame
    tf.word_wrap = True
    
    for i, bullet in enumerate(bullets):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = bullet
        p.font.size = Pt(20)
        p.font.color.rgb = DARK_GRAY
        p.level = 0
        p.space_after = Pt(12)
    
    return slide

def add_section_slide(prs, section_title):
    blank_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(blank_layout)
    
    # Yellow background
    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height)
    bg.fill.solid()
    bg.fill.fore_color.rgb = PADDY_YELLOW
    bg.line.fill.background()
    
    # Green accent bar
    accent = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, Inches(3.5), prs.slide_width, Inches(0.5))
    accent.fill.solid()
    accent.fill.fore_color.rgb = PADDY_GREEN
    accent.line.fill.background()
    
    # Section title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(2.8), Inches(12.333), Inches(1))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = section_title
    p.font.size = Pt(48)
    p.font.bold = True
    p.font.color.rgb = PADDY_GREEN
    p.alignment = PP_ALIGN.CENTER
    
    return slide

# SLIDE 1: Title
add_title_slide(prs, 'Paddy Power: IMC Strategy Analysis', 'Contemporary Marketing Communications | B6BU118')

# SLIDE 2: Introduction
add_content_slide(prs, 'Introduction', [
    'Paddy Power: Irish sports betting company founded in 1988',
    'Known for humorous, controversial marketing campaigns',
    'IMC Strategy: Consistent, bold, cheeky "lad culture" persona',
    'Positions brand as entertainment provider, not just bookmaker',
    '2026 evolution: Balancing entertainment with responsible gambling'
])

# SLIDE 3: Core IMC Pillars
add_content_slide(prs, 'Core IMC Pillars', [
    'Unified Brand Voice: Cheeky, provocative persona across all channels',
    'Omni-channel Integration: Seamless connection between shops, apps, social media',
    'Reactive Marketing: Real-time response to cultural and sporting moments',
    'Brand Positioning: "Mischief" personality that stands out in saturated market'
])

# SLIDE 4: Section - Advertising
add_section_slide(prs, 'Advertising')

# SLIDE 5: Come Out and Play Campaign
add_content_slide(prs, '"Come Out and Play" Campaign (2025)', [
    'Launch: October 2025, created by BBH agency',
    'Cast: Danny Dyer, Coleen Rooney, Gemma Collins, Peter Crouch',
    'Concept: "Queen Vic meets Las Vegas" - dreamlike casino scenario',
    'Danny Dyer as "The House" - cockney casino guide',
    'Channels: TV, Video on Demand, Social Media',
    'Tone: Self-deprecating humor, memorable entertainment'
])

# SLIDE 6: Section - PR & Social Media
add_section_slide(prs, 'Public Relations & Social Media')

# SLIDE 7: PR Stunts
add_content_slide(prs, 'PR Stunts & Earned Media', [
    'Hollywood Sign Stunt (2010 Ryder Cup): Giant hillside sign at Celtic Manor',
    'Result: Millions saw brand without traditional advertising costs',
    'Even Bigger 180 (2024-2026): Linked to PDC World Darts Championship',
    'Promise: £1,000 donation per "180" score to Prostate Cancer UK',
    'Results: Over £1 million raised, thousands of health checks encouraged',
    'Integration: QR codes, newspaper takeovers, charitable giving'
])

# SLIDE 8: Social Media
add_content_slide(prs, 'Social Media Strategy', [
    'Platforms: Instagram, X (Twitter), TikTok for real-time engagement',
    'Content: React to live sports, humorous memes, fan engagement',
    'Example: Katie Taylor boxing fight - video warning about toxic comments',
    'Approach: Tap into sports culture conversations instantly',
    'Result: Highly shareable content that reinforces brand personality'
])

# SLIDE 9: Section - Direct Marketing
add_section_slide(prs, 'Direct Marketing & Sales Promotion')

# SLIDE 10: Direct Marketing
add_content_slide(prs, 'Direct Marketing Channels', [
    'Email: Personalised betting offers and updates',
    'SMS: Targeted promotions and event reminders',
    'Push Notifications: In-app alerts for live events and offers',
    'Purpose: Maintain customer relationships, encourage repeat usage',
    'Integration: Consistent messaging across all direct channels'
])

# SLIDE 11: Paddy\'s Rewards Club
add_content_slide(prs, 'Paddy\'s Rewards Club (Loyalty Program)', [
    'Qualification: Place 5+ bets of £5/€5+ (odds 1/2+) Monday-Sunday',
    'Rewards: Free bet (£5-£50) OR Power Up token (weekly)',
    'Power Up Tokens: Boost odds on eligible bets',
    'Cross-channel: Online, mobile, phone, text, AND shop bets',
    '2025 Update: Now requires opt-in for engaged participation',
    'New Customers: Bet £5, get £30 in free bets (jumps season)'
])

# SLIDE 12: Section - Sponsorship
add_section_slide(prs, 'Sponsorship & Integration')

# SLIDE 13: Sponsorship
add_content_slide(prs, 'Sponsorship Strategy', [
    'PDC World Darts Championship: "Even Bigger 180" campaign integration',
    'First Dates Sponsorship: Extended into 2026 for cultural relevance',
    'Strategy: Link sponsorships with charitable causes and activation',
    'Benefit: Build brand warmth beyond sports betting',
    'Integration: Combine with QR codes, media coverage, social amplification'
])

# SLIDE 14: Integrated Approach
add_content_slide(prs, 'The Integrated Approach', [
    'Pinball Effect: Consumers enter at unpredictable touchpoints',
    'Campaign Flow: Stunt → PR coverage → Social amplification → Direct conversion',
    'Consistent Messaging: Same cheeky personality across all channels',
    'Example: Even Bigger 180 combined sponsorship, charity, QR, PR, social',
    'Result: Unified brand experience regardless of entry point'
])

# SLIDE 15: Responsible Gambling
add_content_slide(prs, 'Responsible Gambling Integration (2026)', [
    'AI-Driven Time Alerts: Remind customers of betting duration',
    'Self-Exclusion Tools: Easy access to responsible gambling controls',
    'Marketing Integration: Responsible messages in all campaigns',
    'Regulatory Alignment: Meeting stricter global regulations',
    'Brand Evolution: Balancing entertainment with corporate responsibility'
])

# SLIDE 16: Conclusion
add_content_slide(prs, 'Conclusion', [
    'Paddy Power demonstrates effective IMC through consistent brand personality',
    'Six IMC tools work together: Advertising, PR, Social, Direct, Promotion, Sponsorship',
    'Key Success: "Mischief" positioning creates distinctive market presence',
    '2026 Evolution: Integrating responsible gambling with entertainment',
    'Lesson: Successful IMC requires channels to work together, not just coexist'
])

# SLIDE 17: Thank You
add_title_slide(prs, 'Thank You', 'Questions & Discussion')

# Save presentation
prs.save('C:/Users/Karen/.openclaw/workspace/Paddy_Power_IMC_Presentation.pptx')
print('PowerPoint presentation created: Paddy_Power_IMC_Presentation.pptx')
print('Total slides: 17')